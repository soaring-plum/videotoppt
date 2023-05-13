import os
import pptx
from pptx.util import Inches

# ppt_filename = input('输入目标ppt文件名(无需后缀)：')

def main(ppt_filename,picPath):
    full_filename = '{}.{}'.format(ppt_filename, 'pptx')
    ppt = pptx.Presentation()#@可以加路径
    for i,j,k in os.walk(picPath):
        pass
        # print(i,k)
    # k=['0.jpg']
    for i in range(len(k)):#新建k张空白页
       ppt.slides.add_slide(ppt.slide_layouts[6])  # 新建空白页
    for i in range(len(k)):
        slide = ppt.slides[i]
        left, top, width, height = pptx.util.Cm(0), pptx.util.Cm(0), pptx.util.Cm(25.4), pptx.util.Cm(19.05)
        slide.shapes.add_picture(image_file=picPath+k[i], left=left, top=top, width=width, height=height)
    ppt.save(full_filename)

if __name__=='__main__' :
    main(ppt_filename='C:/Users/li/Desktop/videotoppt/workspace/pptx/花卉学new',picPath='C:/Users/li/Desktop/videotoppt/workspace/temp2/')